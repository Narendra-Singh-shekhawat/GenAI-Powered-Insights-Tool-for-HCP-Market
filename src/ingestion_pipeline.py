import os
import re
from typing import List, Dict

from sentence_transformers import SentenceTransformer
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.document_loaders import PyPDFLoader
from pptx import Presentation

from weaviate_client import WeaviateClientManager



class IngestionPipeline:
    def __init__(self, weaviate_mode="local"):
        # Embedding model
        self.model = SentenceTransformer("all-MiniLM-L6-v2")

        # Text splitter
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=300,
            chunk_overlap=50
        )

        # Weaviate client
        self.client = WeaviateClientManager(mode=weaviate_mode).get_client()

        # Class name
        self.class_name = "MRInsights"

    # -------------------------------
    # 🔹 SCHEMA SETUP
    # -------------------------------
    def create_schema(self):
        if self.client.schema.exists(self.class_name):
            print("⚠️ Schema already exists")
            return

        schema = {
            "class": self.class_name,
            "vectorizer": "none",
            "properties": [
                {"name": "content", "dataType": ["text"]},
                {"name": "doctor_id", "dataType": ["text"]},
                {"name": "doctor_name", "dataType": ["text"]},
                {"name": "specialty", "dataType": ["text"]},
                {"name": "region", "dataType": ["text"]},
                {"name": "date", "dataType": ["text"]},
                {"name": "product", "dataType": ["text"]},
                {"name": "source", "dataType": ["text"]},
                {"name": "file_name", "dataType": ["text"]}
            ]
        }

        self.client.schema.create_class(schema)
        print("✅ Schema created")

    # -------------------------------
    # 🔹 TEXT CLEANING
    # -------------------------------
    def clean_text(self, text: str) -> str:
        return re.sub(r"\s+", " ", text).strip()

    # -------------------------------
    # 🔹 TRANSCRIPT PARSER
    # -------------------------------
    def parse_transcript(self, text: str) -> Dict:
        def extract(pattern):
            match = re.search(pattern, text)
            return match.group(1).strip() if match else "Unknown"

        metadata = {
            "doctor_id": extract(r"Doctor_ID:\s*(.*)"),
            "doctor_name": extract(r"Doctor:\s*(.*)"),
            "specialty": extract(r"Specialty:\s*(.*)"),
            "region": extract(r"Region:\s*(.*)"),
            "date": extract(r"Date:\s*(.*)"),
            "product": extract(r"Product:\s*(.*)"),
            "source": "transcript",
            "file_name": "N/A"
        }

        content = text.split("Transcript:")[-1]
        content = self.clean_text(content)

        return {"content": content, "metadata": metadata}

    # -------------------------------
    # 🔹 PDF LOADER
    # -------------------------------
    def load_pdf(self, file_path: str) -> List[Dict]:
        loader = PyPDFLoader(file_path)
        docs = loader.load()

        results = []
        for doc in docs:
            results.append({
                "content": self.clean_text(doc.page_content),
                "metadata": {
                    "doctor_id": "Unknown",
                    "doctor_name": "Unknown",
                    "specialty": "Unknown",
                    "region": "Unknown",
                    "date": "Unknown",
                    "product": "Unknown",
                    "source": "pdf",
                    "file_name": os.path.basename(file_path)
                }
            })

        return results

    # -------------------------------
    # 🔹 PPT LOADER
    # -------------------------------
    def load_ppt(self, file_path: str) -> List[Dict]:
        prs = Presentation(file_path)

        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)

        return [{
            "content": self.clean_text(" ".join(text)),
            "metadata": {
                "doctor_id": "Unknown",
                "doctor_name": "Unknown",
                "specialty": "Unknown",
                "region": "Unknown",
                "date": "Unknown",
                "product": "Unknown",
                "source": "ppt",
                "file_name": os.path.basename(file_path)
            }
        }]

    # -------------------------------
    # 🔹 CHUNKING
    # -------------------------------
    def chunk_documents(self, docs: List[Dict]) -> List[Dict]:
        chunks = []

        for doc in docs:
            split_texts = self.splitter.split_text(doc["content"])

            for chunk in split_texts:
                chunks.append({
                    "content": chunk,
                    "metadata": doc["metadata"]
                })

        return chunks

    # -------------------------------
    # 🔹 EMBEDDING
    # -------------------------------
    def embed(self, text: str):
        return self.model.encode(text).tolist()

    # -------------------------------
    # 🔹 STORE IN WEAVIATE
    # -------------------------------
    def store_chunks(self, chunks: List[Dict]):
        for chunk in chunks:
            vector = self.embed(chunk["content"])

            data_object = {
                "content": chunk["content"],
                **chunk["metadata"]
            }

            self.client.data_object.create(
                data_object=data_object,
                class_name=self.class_name,
                vector=vector
            )

        print(f"✅ Stored {len(chunks)} chunks in Weaviate")

    # -------------------------------
    # 🔹 MAIN INGESTION FUNCTION
    # -------------------------------
    def run(self, base_path="data"):
        self.create_schema()

        all_docs = []

        # Transcripts
        transcript_path = os.path.join(base_path, "transcripts")
        if os.path.exists(transcript_path):
            for file in os.listdir(transcript_path):
                if file.endswith(".txt"):
                    with open(os.path.join(transcript_path, file), "r", encoding="utf-8") as f:
                        all_docs.append(self.parse_transcript(f.read()))

        # PDFs
        pdf_path = os.path.join(base_path, "pdfs")
        if os.path.exists(pdf_path):
            for file in os.listdir(pdf_path):
                if file.endswith(".pdf"):
                    all_docs.extend(self.load_pdf(os.path.join(pdf_path, file)))

        # PPTs
        ppt_path = os.path.join(base_path, "ppt")
        if os.path.exists(ppt_path):
            for file in os.listdir(ppt_path):
                if file.endswith(".pptx"):
                    all_docs.extend(self.load_ppt(os.path.join(ppt_path, file)))

        # Chunking
        chunks = self.chunk_documents(all_docs)

        # Store
        self.store_chunks(chunks)


# -------------------------------
# 🚀 RUN SCRIPT
# -------------------------------
if __name__ == "__main__":
    pipeline = IngestionPipeline(weaviate_mode="cloud")  # or "cloud"
    pipeline.run("data")