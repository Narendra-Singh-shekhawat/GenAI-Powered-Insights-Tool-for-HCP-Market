import os
import re
from typing import List, Dict

from sentence_transformers import SentenceTransformer
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Optional loaders
from langchain.document_loaders import PyPDFLoader
from pptx import Presentation


class EmbeddingPipeline:
    def __init__(self):
        self.model = SentenceTransformer("all-MiniLM-L6-v2")
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=300,
            chunk_overlap=50
        )

    # -------------------------------
    # 🔹 TEXT CLEANING
    # -------------------------------
    def clean_text(self, text: str) -> str:
        text = re.sub(r"\s+", " ", text)
        return text.strip()

    # -------------------------------
    # 🔹 TRANSCRIPT PARSER
    # -------------------------------
    def parse_transcript(self, text: str) -> Dict:
        metadata = {}

        def extract(pattern):
            match = re.search(pattern, text)
            return match.group(1).strip() if match else "Unknown"

        metadata["doctor_name"] = extract(r"Doctor: (.*)")
        metadata["specialty"] = extract(r"Specialty: (.*)")
        metadata["region"] = extract(r"Region: (.*)")
        metadata["date"] = extract(r"Date: (.*)")
        metadata["product"] = extract(r"Product: (.*)")
        metadata["doctor_id"] = metadata["doctor_name"].replace(" ", "_")
        metadata["source"] = "transcript"

        content = text.split("Transcript:")[-1]
        content = self.clean_text(content)

        return {"content": content, "metadata": metadata}

    # -------------------------------
    # 🔹 PDF LOADER
    # -------------------------------
    def load_pdf(self, file_path: str) -> List[Dict]:
        loader = PyPDFLoader(file_path)
        docs = loader.load()

        results = []
        for doc in docs:
            content = self.clean_text(doc.page_content)

            metadata = {
                "source": "pdf",
                "file_name": os.path.basename(file_path)
            }

            results.append({"content": content, "metadata": metadata})

        return results

    # -------------------------------
    # 🔹 PPT LOADER
    # -------------------------------
    def load_ppt(self, file_path: str) -> List[Dict]:
        prs = Presentation(file_path)

        slides_text = []
        for slide in prs.slides:
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            slides_text.append(" ".join(text))

        full_text = self.clean_text(" ".join(slides_text))

        return [{
            "content": full_text,
            "metadata": {
                "source": "ppt",
                "file_name": os.path.basename(file_path)
            }
        }]

    # -------------------------------
    # 🔹 CHUNKING
    # -------------------------------
    def chunk_documents(self, docs: List[Dict]) -> List[Dict]:
        chunks = []

        for doc in docs:
            split_texts = self.splitter.split_text(doc["content"])

            for chunk in split_texts:
                chunks.append({
                    "content": chunk,
                    "metadata": doc["metadata"]
                })

        return chunks

    # -------------------------------
    # 🔹 EMBEDDINGS
    # -------------------------------
    def generate_embedding(self, text: str) -> List[float]:
        return self.model.encode(text).tolist()

    # -------------------------------
    # 🔹 MAIN PIPELINE
    # -------------------------------
    def process_transcripts(self, folder_path: str) -> List[Dict]:
        docs = []

        for file in os.listdir(folder_path):
            if file.endswith(".txt"):
                with open(os.path.join(folder_path, file), "r", encoding="utf-8") as f:
                    parsed = self.parse_transcript(f.read())
                    docs.append(parsed)

        return self.chunk_documents(docs)

    def process_pdfs(self, folder_path: str) -> List[Dict]:
        docs = []

        for file in os.listdir(folder_path):
            if file.endswith(".pdf"):
                docs.extend(self.load_pdf(os.path.join(folder_path, file)))

        return self.chunk_documents(docs)

    def process_ppts(self, folder_path: str) -> List[Dict]:
        docs = []

        for file in os.listdir(folder_path):
            if file.endswith(".pptx"):
                docs.extend(self.load_ppt(os.path.join(folder_path, file)))

        return self.chunk_documents(docs)